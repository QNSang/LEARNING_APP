import fitz
from pptx import Presentation
from docx import Document
import os
import re
from typing import List, Dict, Any
import structlog

logger = structlog.get_logger()

# Patterns nhận diện header/footer rác
_NOISE_PATTERNS =[
    r"^trang\s*\d+$",           # "Trang 45"
    r"^\d+\s*/\s*\d+$",         # "45 / 200"  
    r"^-\s*\d+\s*-$",           # "- 45 -"
    r"^\d+$",                   # chỉ là số trang (có thể cần cẩn thận nếu bảng có số nguyên)
]
_NOISE_RE = re.compile("|".join(_NOISE_PATTERNS), re.IGNORECASE)

# Từ không có giá trị tri thức
_SKIP_PHRASES = {
    "xem thêm", "tài liệu tham khảo", "nguồn:", "references",
    "continued", "tiếp theo", "..."
}

def _is_noise(text: str) -> bool:
    t = text.strip()
    if not t:
        return True
    if _NOISE_RE.match(t):
        return True
    if t.lower() in _SKIP_PHRASES:
        return True
    return False

def _detect_repeating_lines(doc: fitz.Document, min_ratio: float = 0.4) -> set:
    """
    Tìm các dòng text xuất hiện ở số trang >= min_ratio * total_pages.

    """
    from collections import Counter
    line_counter = Counter()
    total_pages = len(doc)
    
    # Sample tối đa 30 trang
    sample_pages = min(30, total_pages)
    if sample_pages == 0:
        return set()

    for page in doc[:sample_pages]:
        for block in page.get_text("blocks"):
            if block[6] == 0:  # text block
                line = block[4].strip()
                if line:
                    line_counter[line] += 1
    
    # Tính threshold động dựa trên tổng số trang (ít nhất xuất hiện 4 lần)
    threshold = max(4, int(sample_pages * min_ratio))
    
    # Chỉ đánh dấu là noise header/footer nếu nó có độ dài từ 5 đến 100 ký tự.
    # (Nếu quá ngắn có thể là tên biến, quá dài có thể là 1 đoạn văn lặp lại hợp lệ)
    return {
        line for line, count in line_counter.items() 
        if count >= threshold and 5 < len(line) < 100
    }

class ParserService:

    # ─────────────────────────────────────────
    # PDF
    # ─────────────────────────────────────────
    @staticmethod
    async def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        pages =[]
        file_name = os.path.basename(file_path)
        
        try:
            doc = fitz.open(file_path)
            # Dùng logic detect noise mới
            repeating_lines = _detect_repeating_lines(doc)
            
            for page_num, page in enumerate(doc):
                blocks = page.get_text("blocks")  
                lines = []
                for block in blocks:
                    if block[6] != 0:   # bỏ image blocks
                        continue
                    text = block[4].strip()
                    
                    if _is_noise(text):
                        continue
                    if text in repeating_lines:
                        continue
                        
                    lines.append(text)
                
                content = "\n".join(lines).strip()
                
                if content and len(content.encode("utf-8")) / len(content) > 3:
                    logger.warning("encoding_issue", page=page_num + 1, file=file_name)
                
                if content:
                    pages.append({
                        "page_num": page_num + 1,
                        "content": content,
                        "metadata": {
                            "source": "pdf",
                            "page": page_num + 1,
                            "file_name": file_name,
                            "source_ref": f"{file_name} — trang {page_num + 1}"
                        }
                    })
            doc.close()
        except Exception as e:
            logger.error("pdf_parsing_failed", error=str(e), file=file_path)
            raise
            
        return pages

    # ─────────────────────────────────────────
    # PPTX  
    # ─────────────────────────────────────────
    @staticmethod
    async def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        slides =[]
        file_name = os.path.basename(file_path)
        
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                title = ""
                bullets =[]
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                            ph_idx = shape.placeholder_format.idx
                            if ph_idx == 0:  # title
                                title = shape.text.strip()
                                continue
                        
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and not _is_noise(text):
                                bullets.append(text)
                    
                    # FIX: Trích xuất nội dung từ Bảng biểu (Rất hay dùng trong Slide học thuật)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_data =[cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_data:
                                bullets.append(" | ".join(row_data))
                
                parts =[]
                if title:
                    parts.append(f"[Slide {i+1}] {title}")
                parts.extend(bullets)
                
                content = "\n".join(parts).strip()
                if content:
                    slides.append({
                        "page_num": i + 1,
                        "content": content,
                        "metadata": {
                            "source": "pptx",
                            "page": i + 1,
                            "file_name": file_name,
                            "source_ref": f"{file_name} — slide {i + 1}: {title or 'không có tiêu đề'}"
                        }
                    })
        except Exception as e:
            logger.error("pptx_parsing_failed", error=str(e), file=file_path)
            raise
            
        return slides

    # ─────────────────────────────────────────
    # DOCX
    # ─────────────────────────────────────────
    @staticmethod
    async def parse_docx(file_path: str) -> List[Dict[str, Any]]:
        sections =[]
        file_name = os.path.basename(file_path)
        
        try:
            doc = Document(file_path)
            current_section =[]
            current_heading = "Mở đầu"
            section_num = 1
            
            # Đọc Paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                    
                is_heading = para.style.name.startswith("Heading")
                
                if is_heading:
                    if current_section:
                        sections.append({
                            "page_num": section_num,
                            "content": "\n".join(current_section).strip(),
                            "metadata": {
                                "source": "docx",
                                "file_name": file_name,
                                "heading": current_heading,
                                "source_ref": f"{file_name} — {current_heading}"
                            }
                        })
                        current_section =[]
                        section_num += 1
                    
                    current_heading = text
                    current_section.append(f"[{text}]")
                else:
                    if not _is_noise(text):
                        current_section.append(text)
            
            # Đẩy nốt phần tử đoạn văn cuối
            if current_section:
                sections.append({
                    "page_num": section_num,
                    "content": "\n".join(current_section).strip(),
                    "metadata": {
                        "source": "docx",
                        "file_name": file_name,
                        "heading": current_heading,
                        "source_ref": f"{file_name} — {current_heading}"
                    }
                })
                
    
            table_section =[]
            for table in doc.tables:
                table_section.append("[Bảng dữ liệu]")
                for row in table.rows:
                    row_data =[cell.text.strip().replace("\n", " ") for cell in row.cells if cell.text.strip()]
                    if row_data:
                        table_section.append(" | ".join(row_data))
                table_section.append("") # Xuống dòng phân tách các bảng
            
            if table_section and "".join(table_section).strip() != "[Bảng dữ liệu]":
                sections.append({
                    "page_num": section_num + 1,
                    "content": "\n".join(table_section).strip(),
                    "metadata": {
                        "source": "docx",
                        "file_name": file_name,
                        "heading": "Các bảng biểu trong tài liệu",
                        "source_ref": f"{file_name} — Bảng biểu"
                    }
                })

        except Exception as e:
            logger.error("docx_parsing_failed", error=str(e), file=file_path)
            raise
            
        return sections

    # ─────────────────────────────────────────
    # TXT (Giữ nguyên)
    # ─────────────────────────────────────────
    @staticmethod
    async def parse_txt(file_path: str) -> List[Dict[str, Any]]:
        file_name = os.path.basename(file_path)
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read().strip()
            
            raw_chunks = re.split(r"\n---+\n|\n\n", content)
            chunks =[c.strip() for c in raw_chunks if c.strip() and not _is_noise(c.strip())]
            
            return[{
                "page_num": i + 1,
                "content": chunk,
                "metadata": {
                    "source": "txt",
                    "file_name": file_name,
                    "source_ref": f"{file_name} — phần {i + 1}"
                }
            } for i, chunk in enumerate(chunks)]
        except Exception as e:
            logger.error("txt_parsing_failed", error=str(e), file=file_path)
            raise

    # ─────────────────────────────────────────
    # Dispatcher
    # ─────────────────────────────────────────
    async def parse_file(self, file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        parsers = {
            ".pdf":  self.parse_pdf,
            ".pptx": self.parse_pptx,
            ".docx": self.parse_docx,
            ".txt":  self.parse_txt,
        }
        if ext not in parsers:
            raise ValueError(f"Unsupported file type: {ext}")
        
        result = await parsers[ext](file_path)
        logger.info("parsing_done", file=os.path.basename(file_path), pages=len(result))
        return result